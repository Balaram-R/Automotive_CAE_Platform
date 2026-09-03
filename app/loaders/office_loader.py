"""
app/loaders/office_loader.py
============================
XLSX / PPTX loaders.
"""

from __future__ import annotations

import os
import tempfile
from pathlib import Path

from app.loaders.base import BaseLoader
from app.loaders.loader_factory import LoaderFactory
from app.loaders.image_loader import _groq_vision_caption
from app.models.schemas import FileMetadata, RawDocument
from app.utils.logging import get_logger

logger = get_logger("loaders.office")


class ExcelLoader(BaseLoader):
    supported_extensions = [".xlsx", ".xls"]

    def load(self, filepath: str, metadata: FileMetadata, **kw) -> list[RawDocument]:
        try:
            import openpyxl
        except ImportError:
            logger.error("pip install openpyxl")
            return []
        try:
            wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
            parts: list[str] = []
            for name in wb.sheetnames:
                rows = [
                    "\t".join(str(c) if c is not None else "" for c in row)
                    for row in wb[name].iter_rows(values_only=True)
                ]
                if rows:
                    parts.append(f"[Sheet: {name}]\n" + "\n".join(rows))
            wb.close()
            return [RawDocument(content="\n\n".join(parts), metadata=metadata, source_loader="ExcelLoader")] if parts else []
        except Exception as exc:
            logger.error("Excel load failed %s: %s", filepath, exc)
            return []


class PPTXLoader(BaseLoader):
    supported_extensions = [".pptx", ".ppt"]

    def _shape_text(self, shape) -> list[str]:
        """Recursively extract text from a shape, including grouped shapes."""
        texts: list[str] = []
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            for sub in shape.shapes:
                texts.extend(self._shape_text(sub))
            return texts
        if shape.has_text_frame:
            texts.extend(p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip())
        if shape.has_table:
            for row in shape.table.rows:
                texts.append("\t".join(c.text.strip() for c in row.cells))
        return texts

    def load(self, filepath: str, metadata: FileMetadata, **kw) -> list[RawDocument]:
        try:
            from pptx import Presentation
        except ImportError:
            logger.error("pip install python-pptx")
            return []
        try:
            prs = Presentation(filepath)
            slides: list[str] = []
            image_only_slides = 0
            # Limit how many image-only slides get vision-captioned (each call
            # costs ~10s). Default: caption up to 12 slides, sampled evenly.
            max_caption = int(kw.get("vision_max_slides", 12))
            total_slides = len(prs.slides._sldIdLst) if hasattr(prs.slides, "_sldIdLst") else len(list(prs.slides))
            caption_step = max(1, total_slides // max_caption) if max_caption > 0 else 0
            caption_count = 0
            for i, slide in enumerate(prs.slides):
                texts: list[str] = []
                has_image = False
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # PICTURE
                        has_image = True
                    texts.extend(self._shape_text(shape))
                # Include speaker notes
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        texts.append(f"[Notes]\n{notes}")
                if texts:
                    slides.append(f"[Slide {i + 1}]\n" + "\n".join(texts))
                elif has_image:
                    image_only_slides += 1
                    # Vision-caption a sampled subset of image-only slides so
                    # they get real content instead of a placeholder.
                    if max_caption > 0 and caption_count < max_caption and (i % caption_step == 0):
                        caption = self._caption_slide_image(slide, i, kw)
                        if caption:
                            slides.append(f"[Slide {i + 1}]\n{caption}")
                            caption_count += 1

            if slides:
                return [RawDocument(content="\n\n".join(slides), metadata=metadata, source_loader="PPTXLoader")]

            # Image-only deck with no vision captions — produce a placeholder.
            if image_only_slides:
                content = (
                    f"[Presentation: {metadata.filename}]\n"
                    f"Image-only slide deck with {image_only_slides} slide(s) containing images "
                    f"but no extractable text. OCR is required to read the slide content."
                )
                return [RawDocument(content=content, metadata=metadata, source_loader="PPTXLoader")]

            return []
        except Exception as exc:
            logger.error("PPTX load failed %s: %s", filepath, exc)
            return []

    def _caption_slide_image(self, slide, slide_index: int, kw: dict) -> str:
        """Caption the first picture on a slide using the Groq vision model."""
        cfg = kw.get("config")
        vision_enabled = True
        vision_model = "qwen/qwen3.6-27b"
        api_key = os.environ.get("GROQ_API_KEY", "")
        if cfg is not None:
            llm = getattr(cfg, "llm", None)
            if llm is not None:
                vision_enabled = getattr(llm, "vision_enabled", True)
                vision_model = getattr(llm, "vision_model", None) or vision_model
                api_key = getattr(llm, "api_key", None) or api_key
        if not vision_enabled or not api_key:
            return ""

        # Find the first picture shape and save its image to a temp file
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                try:
                    img = shape.image
                    ext = img.content_type.split("/")[-1] or "png"
                    if ext == "jpeg":
                        ext = "jpg"
                    with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as f:
                        f.write(img.blob)
                        tmp = f.name
                    try:
                        caption = _groq_vision_caption(tmp, vision_model, api_key, max_tokens=400)
                    finally:
                        try:
                            os.unlink(tmp)
                        except OSError:
                            pass
                    if caption:
                        return f"[Vision Description]\n{caption}"
                except Exception as exc:
                    logger.warning("Vision caption failed for slide %d: %s", slide_index + 1, exc)
                break
        return ""


LoaderFactory.register_many(ExcelLoader.supported_extensions, ExcelLoader())
LoaderFactory.register_many(PPTXLoader.supported_extensions, PPTXLoader())
